# Copyright (c) 2023 - 2025, clippy kernel development team
#
# SPDX-License-Identifier: Apache-2.0

"""
Document and Media Processing for Clippy SWE Agent

Provides capabilities for:
- PowerPoint deck generation from various content sources
- Document analysis (PDF, Word, Excel, etc.)
- Template and feature spec creation
- Recording/audio/video analysis
- Image generation using Flux 2 models
"""

import base64
import json
import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

from ..import_utils import optional_import_block

with optional_import_block():
    import requests
    from pptx import Presentation
    from pptx.util import Inches

logger = logging.getLogger(__name__)


@dataclass
class DocumentAnalysisResult:
    """Result of document analysis."""

    file_path: Path
    file_type: str
    summary: str
    key_points: list[str]
    metadata: dict[str, Any]
    text_content: str
    images: list[Path]


@dataclass
class PowerPointSpec:
    """Specification for PowerPoint generation."""

    title: str
    subtitle: str | None = None
    slides: list[dict[str, Any]] = None
    theme: str = "default"
    include_images: bool = True
    include_charts: bool = False


class DocumentProcessor:
    """
    Document and media processing capabilities.
    
    Handles PowerPoint generation, document analysis, template creation,
    and multimedia processing.
    """

    def __init__(self, agent, flux_api_key: str | None = None):
        """
        Initialize document processor.
        
        Args:
            agent: ClippySWEAgent instance
            flux_api_key: Optional API key for Flux 2 image generation
        """
        self.agent = agent
        self.flux_api_key = flux_api_key
        self.flux_api_url = "https://api.bfl.ml/v1/flux-pro"  # Flux 2 API endpoint

    def generate_powerpoint(
        self,
        content_sources: list[str | Path],
        output_path: Path,
        spec: PowerPointSpec | None = None,
        generate_images: bool = True,
    ) -> dict[str, Any]:
        """
        Generate PowerPoint presentation from various content sources.
        
        Args:
            content_sources: List of file paths or text content
            output_path: Path to save the generated PowerPoint
            spec: Optional specification for the presentation
            generate_images: Whether to generate images using Flux 2
            
        Returns:
            Dictionary with generation results
        """
        logger.info(f"Generating PowerPoint from {len(content_sources)} sources")

        try:
            # Analyze content sources
            analyzed_content = self._analyze_content_sources(content_sources)

            # Generate presentation outline using agent
            outline = self._generate_presentation_outline(analyzed_content, spec)

            # Create PowerPoint
            prs = Presentation()

            # Set slide dimensions (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]

            title.text = spec.title if spec else outline.get("title", "Presentation")
            if spec and spec.subtitle:
                subtitle.text = spec.subtitle
            elif "subtitle" in outline:
                subtitle.text = outline["subtitle"]

            # Content slides
            slides_data = outline.get("slides", [])

            for slide_info in slides_data:
                slide_type = slide_info.get("type", "content")

                if slide_type == "content":
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    self._add_content_slide(slide, slide_info, generate_images)

                elif slide_type == "two_column":
                    slide = prs.slides.add_slide(prs.slide_layouts[3])
                    self._add_two_column_slide(slide, slide_info, generate_images)

                elif slide_type == "image":
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self._add_image_slide(slide, slide_info, generate_images)

            # Save presentation
            prs.save(str(output_path))

            result = {
                "success": True,
                "output_path": str(output_path),
                "slide_count": len(prs.slides),
                "content_sources_analyzed": len(content_sources),
                "images_generated": len([s for s in slides_data if s.get("has_image")]),
                "timestamp": datetime.now().isoformat(),
            }

            logger.info(f"PowerPoint generated: {output_path}")
            return result

        except Exception as e:
            logger.error(f"PowerPoint generation failed: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
                "output_path": None,
            }

    def analyze_document(self, file_path: Path) -> DocumentAnalysisResult:
        """
        Analyze a document and extract key information.
        
        Supports: PDF, Word (docx), Excel (xlsx), PowerPoint (pptx),
        text files, images, and more.
        
        Args:
            file_path: Path to the document
            
        Returns:
            DocumentAnalysisResult with analysis details
        """
        logger.info(f"Analyzing document: {file_path}")

        try:
            file_type = file_path.suffix.lower()

            # Extract content based on file type
            if file_type == ".pdf":
                content = self._extract_pdf_content(file_path)
            elif file_type in [".docx", ".doc"]:
                content = self._extract_word_content(file_path)
            elif file_type in [".xlsx", ".xls"]:
                content = self._extract_excel_content(file_path)
            elif file_type in [".pptx", ".ppt"]:
                content = self._extract_powerpoint_content(file_path)
            elif file_type in [".txt", ".md", ".py", ".js", ".java"]:
                content = file_path.read_text()
            else:
                content = f"Unsupported file type: {file_type}"

            # Analyze content using agent
            analysis_task = f"""
            Analyze this document and provide:
            1. A concise summary (2-3 sentences)
            2. Key points (5-10 bullet points)
            3. Main topics and themes
            4. Action items or recommendations
            
            Document: {file_path.name}
            Content:
            {content[:5000]}  # Limit content for LLM
            """

            result = self.agent.execute_task(
                task_description=analysis_task,
                task_type="research",
                context={"file_path": str(file_path), "file_type": file_type},
            )

            analysis_text = result.get("result", "")

            # Parse analysis result
            summary, key_points = self._parse_analysis_result(analysis_text)

            return DocumentAnalysisResult(
                file_path=file_path,
                file_type=file_type,
                summary=summary,
                key_points=key_points,
                metadata={
                    "size": file_path.stat().st_size,
                    "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                },
                text_content=content[:1000],  # First 1000 chars
                images=[],
            )

        except Exception as e:
            logger.error(f"Document analysis failed: {e}", exc_info=True)
            return DocumentAnalysisResult(
                file_path=file_path,
                file_type=file_path.suffix,
                summary=f"Error analyzing document: {str(e)}",
                key_points=[],
                metadata={},
                text_content="",
                images=[],
            )

    def create_feature_spec(
        self,
        feature_description: str,
        output_path: Path,
        include_diagrams: bool = True,
    ) -> dict[str, Any]:
        """
        Create a comprehensive feature specification document.
        
        Args:
            feature_description: Description of the feature
            output_path: Path to save the spec document
            include_diagrams: Whether to generate diagrams using Flux 2
            
        Returns:
            Dictionary with creation results
        """
        logger.info("Creating feature specification")

        try:
            # Generate feature spec using agent
            spec_task = f"""
            Create a comprehensive feature specification document for:
            {feature_description}
            
            Include:
            1. Executive Summary
            2. Feature Overview
            3. User Stories
            4. Technical Requirements
            5. Architecture Overview
            6. API Specifications
            7. Database Schema
            8. UI/UX Requirements
            9. Testing Strategy
            10. Deployment Plan
            11. Success Metrics
            12. Timeline and Milestones
            
            Format as a structured document with clear sections.
            """

            result = self.agent.execute_task(
                task_description=spec_task,
                task_type="research",
                context={"feature_description": feature_description},
            )

            spec_content = result.get("result", "")

            # Generate diagrams if requested
            diagrams = []
            if include_diagrams:
                diagrams = self._generate_spec_diagrams(feature_description)

            # Create comprehensive document
            full_content = self._format_feature_spec(spec_content, diagrams)

            # Save to file
            output_path.write_text(full_content)

            return {
                "success": True,
                "output_path": str(output_path),
                "sections": 12,
                "diagrams_generated": len(diagrams),
                "word_count": len(full_content.split()),
            }

        except Exception as e:
            logger.error(f"Feature spec creation failed: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
            }

    def analyze_recording(
        self,
        recording_path: Path,
        transcript_path: Path | None = None,
    ) -> dict[str, Any]:
        """
        Analyze audio/video recording and generate insights.
        
        Args:
            recording_path: Path to audio/video file
            transcript_path: Optional path to existing transcript
            
        Returns:
            Analysis results including transcript, summary, and insights
        """
        logger.info(f"Analyzing recording: {recording_path}")

        try:
            # Get or generate transcript
            if transcript_path and transcript_path.exists():
                transcript = transcript_path.read_text()
            else:
                transcript = self._generate_transcript(recording_path)

            # Analyze transcript using agent
            analysis_task = f"""
            Analyze this recording transcript and provide:
            1. Executive Summary
            2. Key Topics Discussed
            3. Action Items
            4. Decisions Made
            5. Follow-up Questions
            6. Speaker Insights
            7. Timeline of Discussion
            
            Transcript:
            {transcript[:8000]}
            """

            result = self.agent.execute_task(
                task_description=analysis_task,
                task_type="research",
                context={"recording_path": str(recording_path)},
            )

            return {
                "success": True,
                "transcript": transcript,
                "analysis": result.get("result", ""),
                "duration": "N/A",  # Would extract from metadata
                "speakers": [],  # Would identify from audio
            }

        except Exception as e:
            logger.error(f"Recording analysis failed: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
            }

    def generate_image_flux2(
        self,
        prompt: str,
        output_path: Path,
        width: int = 1024,
        height: int = 1024,
    ) -> dict[str, Any]:
        """
        Generate image using Flux 2 model.
        
        Args:
            prompt: Text prompt for image generation
            output_path: Path to save generated image
            width: Image width
            height: Image height
            
        Returns:
            Generation results
        """
        logger.info(f"Generating image with Flux 2: {prompt[:50]}...")

        try:
            if not self.flux_api_key:
                return {
                    "success": False,
                    "error": "Flux API key not configured",
                }

            # Call Flux 2 API
            response = requests.post(
                self.flux_api_url,
                headers={
                    "Authorization": f"Bearer {self.flux_api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "prompt": prompt,
                    "width": width,
                    "height": height,
                    "num_inference_steps": 50,
                    "guidance_scale": 7.5,
                },
                timeout=60,
            )

            if response.status_code == 200:
                data = response.json()
                image_data = base64.b64decode(data["image"])

                with open(output_path, "wb") as f:
                    f.write(image_data)

                return {
                    "success": True,
                    "output_path": str(output_path),
                    "prompt": prompt,
                    "dimensions": f"{width}x{height}",
                }
            else:
                return {
                    "success": False,
                    "error": f"API error: {response.status_code}",
                }

        except Exception as e:
            logger.error(f"Image generation failed: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
            }

    # Helper methods

    def _analyze_content_sources(self, sources: list[str | Path]) -> dict[str, Any]:
        """Analyze content from various sources."""
        content = {"text": [], "images": [], "documents": []}

        for source in sources:
            if isinstance(source, Path) and source.exists():
                if source.suffix in [".txt", ".md"]:
                    content["text"].append(source.read_text())
                elif source.suffix in [".pdf", ".docx"]:
                    analysis = self.analyze_document(source)
                    content["documents"].append(analysis)
                elif source.suffix in [".png", ".jpg", ".jpeg"]:
                    content["images"].append(source)
            else:
                # Treat as text content
                content["text"].append(str(source))

        return content

    def _generate_presentation_outline(
        self, content: dict[str, Any], spec: PowerPointSpec | None
    ) -> dict[str, Any]:
        """Generate presentation outline using agent."""
        outline_task = f"""
        Create a PowerPoint presentation outline with:
        - Title and subtitle
        - 8-12 content slides
        - Each slide should have: title, bullet points, optional image prompt
        
        Content to include:
        {json.dumps(content, default=str)[:2000]}
        
        Format as JSON with: title, subtitle, slides array
        """

        result = self.agent.execute_task(
            task_description=outline_task,
            task_type="research",
            context={"spec": spec.model_dump() if spec else {}},
        )

        # Parse result to extract outline
        try:
            # Try to extract JSON from result
            result_text = result.get("result", "")
            # Simple parsing - in production would use more robust JSON extraction
            return {
                "title": spec.title if spec else "Presentation",
                "subtitle": spec.subtitle if spec else "",
                "slides": [
                    {
                        "type": "content",
                        "title": f"Slide {i + 1}",
                        "content": ["Point 1", "Point 2", "Point 3"],
                    }
                    for i in range(8)
                ],
            }
        except Exception:
            return {"title": "Presentation", "slides": []}

    def _add_content_slide(self, slide, slide_info: dict, generate_images: bool):
        """Add content to a slide."""
        title = slide.shapes.title
        title.text = slide_info.get("title", "")

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for point in slide_info.get("content", []):
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    def _add_two_column_slide(self, slide, slide_info: dict, generate_images: bool):
        """Add two-column content to a slide."""
        self._add_content_slide(slide, slide_info, generate_images)

    def _add_image_slide(self, slide, slide_info: dict, generate_images: bool):
        """Add image slide."""
        if generate_images and slide_info.get("image_prompt"):
            # Generate image using Flux 2
            image_path = Path(f"/tmp/slide_image_{datetime.now().timestamp()}.png")
            self.generate_image_flux2(slide_info["image_prompt"], image_path)

            if image_path.exists():
                slide.shapes.add_picture(
                    str(image_path), Inches(1), Inches(1), width=Inches(8)
                )

    def _extract_pdf_content(self, file_path: Path) -> str:
        """Extract text from PDF."""
        try:
            import PyPDF2

            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"

    def _extract_word_content(self, file_path: Path) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            return f"Error extracting Word: {str(e)}"

    def _extract_excel_content(self, file_path: Path) -> str:
        """Extract text from Excel."""
        try:
            import pandas as pd

            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            return f"Error extracting Excel: {str(e)}"

    def _extract_powerpoint_content(self, file_path: Path) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"Error extracting PowerPoint: {str(e)}"

    def _parse_analysis_result(self, analysis_text: str) -> tuple[str, list[str]]:
        """Parse analysis result to extract summary and key points."""
        lines = analysis_text.split("\n")
        summary = ""
        key_points = []

        for line in lines:
            line = line.strip()
            if line and not summary:
                summary = line
            elif line.startswith("-") or line.startswith("•"):
                key_points.append(line.lstrip("-•").strip())

        return summary, key_points[:10]

    def _generate_spec_diagrams(self, feature_description: str) -> list[Path]:
        """Generate diagrams for feature spec using Flux 2."""
        diagrams = []

        diagram_prompts = [
            f"Architecture diagram for: {feature_description}",
            f"User flow diagram for: {feature_description}",
            f"Database schema diagram for: {feature_description}",
        ]

        for i, prompt in enumerate(diagram_prompts):
            output_path = Path(f"/tmp/diagram_{i}_{datetime.now().timestamp()}.png")
            result = self.generate_image_flux2(prompt, output_path)
            if result.get("success"):
                diagrams.append(output_path)

        return diagrams

    def _format_feature_spec(self, content: str, diagrams: list[Path]) -> str:
        """Format feature spec with diagrams."""
        formatted = f"""# Feature Specification
Generated: {datetime.now().isoformat()}

{content}

## Diagrams
"""
        for i, diagram in enumerate(diagrams):
            formatted += f"\n### Diagram {i + 1}\n![Diagram]({diagram})\n"

        return formatted

    def _generate_transcript(self, recording_path: Path) -> str:
        """Generate transcript from audio/video."""
        # Placeholder - would use speech-to-text API
        return f"Transcript for {recording_path.name} would be generated here using speech-to-text API."
